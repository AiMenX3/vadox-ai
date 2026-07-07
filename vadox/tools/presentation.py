"""
Professionelle Präsentationen mit python-pptx.
Dunkles Theme mit weißem Text — explizit im Theme verankert.
"""
import os
import json
import zipfile
import shutil
import tempfile
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Farben ────────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)
C_BG2       = RGBColor(0x1A, 0x28, 0x3C)
C_ACCENT    = RGBColor(0x00, 0xD4, 0xFF)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_OFFWHITE  = RGBColor(0xE0, 0xF0, 0xFF)
C_SUBTEXT   = RGBColor(0x90, 0xC0, 0xE0)
C_GOLD      = RGBColor(0xFF, 0xC8, 0x50)
C_LINE      = RGBColor(0x00, 0x60, 0x88)

W, H = 13.33, 7.5


# ── Hilfsfunktionen ──────────────────────────────────────────────────────────

def _patch_theme(pptx_path: str):
    """
    Ändert dk1 (Standard-Textfarbe) im Theme direkt auf Weiß.
    Dadurch kann PowerPoint den Text NICHT mehr auf Schwarz setzen.
    """
    tmp = pptx_path + ".tmp"
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'ppt/theme/theme1.xml':
                    # dk1 auf Weiß setzen (Standard-Textfarbe)
                    xml = data.decode('utf-8')
                    xml = xml.replace(
                        '<a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>',
                        '<a:dk1><a:srgbClr val="FFFFFF"/></a:dk1>'
                    )
                    # lt1 auf Dunkelblau setzen (Standard-Hintergrundfarbe)
                    xml = xml.replace(
                        '<a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1>',
                        '<a:lt1><a:srgbClr val="0D1B2A"/></a:lt1>'
                    )
                    data = xml.encode('utf-8')
                zout.writestr(item, data)
    os.replace(tmp, pptx_path)


def _solid_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color: RGBColor):
    s = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _text(slide, text, left, top, width, height,
          size=18, bold=False, color=C_WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = align

    r = p.add_run()
    r.text = str(text)
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.name      = "Calibri"
    r.font.color.rgb = color
    return tb


def _add_image(slide, img_path, left, top, width, height, behind=False):
    if not img_path or not os.path.exists(img_path):
        return False
    try:
        pic = slide.shapes.add_picture(
            img_path,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        if behind:
            sp = slide.shapes._spTree
            sp.remove(pic._element)
            sp.insert(2, pic._element)
        return True
    except Exception:
        return False


# ── Folien ────────────────────────────────────────────────────────────────────

def _slide_title(prs, title, subtitle, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)

    # Bild als Hintergrund rechts
    _add_image(slide, img_path, W*0.45, 0, W*0.55, H, behind=True)

    # Dunkle linke Hälfte (Textbereich)
    _rect(slide, 0, 0, W*0.55, H, C_BG)

    # Akzentstreifen links
    _rect(slide, 0, 0, 0.22, H, C_ACCENT)

    # Linie oben
    _rect(slide, 0, 0, W, 0.07, C_ACCENT)

    # Titel
    _text(slide, title, 0.5, 1.8, W*0.52, 2.0,
          size=44, bold=True, color=C_WHITE)

    # Trennlinie
    _rect(slide, 0.5, 4.0, 5.5, 0.05, C_ACCENT)

    # Untertitel
    _text(slide, subtitle, 0.5, 4.2, W*0.52, 1.0,
          size=20, color=C_OFFWHITE)

    # Datum
    _text(slide, datetime.now().strftime("%d. %B %Y"),
          0.5, H-0.65, 4.5, 0.45, size=11, color=C_SUBTEXT)

    # VADOX
    _text(slide, "VADOX", W*0.55-2.8, H-0.65, 2.5, 0.45,
          size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.RIGHT)
    return slide


def _slide_content(prs, title, bullets, img_path, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)

    has_img = img_path and os.path.exists(img_path)
    txt_w   = W * 0.54 if has_img else W - 0.8

    if has_img:
        _add_image(slide, img_path, W*0.57, 0, W*0.43, H, behind=True)
        # Abdunkeln des Bildes
        s = slide.shapes.add_shape(
            1, Inches(W*0.57), Inches(0), Inches(W*0.43), Inches(H)
        )
        s.fill.solid()
        s.fill.fore_color.rgb = C_BG
        # 40% Deckkraft über XML
        from lxml import etree
        from pptx.oxml.ns import qn
        spPr = s._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if spPr is not None:
            srgb = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                ae = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                ae.set('val', '40000')
        s.line.fill.background()
        _rect(slide, W*0.57, 0, 0.03, H, C_LINE)

    # Kopfzeile
    _rect(slide, 0, 0, W, 1.08, C_BG2)
    _rect(slide, 0, 0, 0.14, 1.08, C_ACCENT)
    _rect(slide, 0, 1.08, txt_w, 0.04, C_ACCENT)

    # Titel in Kopfzeile — WEISS
    _text(slide, title, 0.28, 0.12, txt_w-0.4, 0.82,
          size=25, bold=True, color=C_WHITE)

    # Seitenzahl
    _text(slide, f"{num}/{total}", txt_w-1.3, 0.25, 1.1, 0.5,
          size=12, color=C_ACCENT, align=PP_ALIGN.RIGHT)

    # Bullets
    y = 1.28
    for i, b in enumerate(bullets[:8]):
        if not b.strip():
            continue
        # Nummernkreis
        circ = slide.shapes.add_shape(
            9, Inches(0.25), Inches(y+0.05), Inches(0.38), Inches(0.38)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = C_ACCENT
        circ.line.fill.background()
        # Nummer — DUNKEL auf CYAN-Kreis
        _text(slide, str(i+1), 0.25, y+0.03, 0.38, 0.38,
              size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        # Bullet-Text — HELL auf DUNKEL
        _text(slide, b, 0.75, y+0.02, txt_w-0.9, 0.58,
              size=15, color=C_OFFWHITE)
        y += 0.67
        if y > H-0.6:
            break

    # Fußzeile
    _rect(slide, 0, H-0.3, W, 0.3, C_BG2)
    _rect(slide, 0, H-0.3, W, 0.025, C_LINE)
    _text(slide, "VADOX KI-Assistent", 0.3, H-0.27, 4.0, 0.24,
          size=9, color=C_LINE)
    return slide


def _slide_closing(prs, main_title, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)
    _add_image(slide, img_path, W*0.45, 0, W*0.55, H, behind=True)
    _rect(slide, 0, 0, W*0.55, H, C_BG)
    _rect(slide, 0, 0, 0.22, H, C_GOLD)
    _rect(slide, 0, 0, W, 0.07, C_GOLD)

    _text(slide, "Vielen Dank!", 0.5, 1.9, W*0.52, 1.8,
          size=52, bold=True, color=C_WHITE)
    _rect(slide, 0.5, 3.85, 5.5, 0.06, C_GOLD)
    _text(slide, main_title, 0.5, 4.05, W*0.52, 0.9,
          size=20, italic=True, color=C_OFFWHITE)
    _text(slide, "Erstellt mit VADOX KI-Assistent",
          0.5, H-0.65, 5.0, 0.42, size=11, color=C_SUBTEXT)
    _text(slide, "VADOX", W*0.55-2.8, H-0.65, 2.5, 0.42,
          size=12, bold=True, color=C_GOLD, align=PP_ALIGN.RIGHT)
    return slide


# ── Haupt-Funktion ────────────────────────────────────────────────────────────

def create_presentation(title: str, slides_data: list, save_path: str = None) -> str:
    from vadox.tools.image_fetcher import fetch_images_for_presentation, cleanup_temp_images
    try:
        prs = Presentation()
        prs.slide_width  = Inches(W)
        prs.slide_height = Inches(H)

        n_imgs = len(slides_data) + 2
        topic  = slides_data[0].get("image_query", title) if slides_data else title
        images = fetch_images_for_presentation(topic, count=n_imgs)
        img_i  = 0

        # Titelfolie
        subtitle = (slides_data[0].get("subtitle", "") if slides_data else "") or "Erstellt von Vadox"
        _slide_title(prs, title, subtitle, images[img_i] if images else None)
        img_i += 1

        # Inhaltsfolien
        content = slides_data[1:] if len(slides_data) > 1 else slides_data
        total   = len(content) + 2
        for i, s in enumerate(content):
            img = images[img_i] if img_i < len(images) else None
            img_i += 1
            _slide_content(prs,
                title   = s.get("title", f"Punkt {i+1}"),
                bullets = s.get("bullets", []),
                img_path= img,
                num=i+2, total=total)

        # Schlussfolie
        _slide_closing(prs, title, images[img_i] if img_i < len(images) else None)

        # Speichern
        if not save_path:
            ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe = "".join(c for c in title if c.isalnum() or c in " _-")[:30].strip()
            save_path = str(Path.home() / "Desktop" / f"{safe}_{ts}.pptx")

        prs.save(save_path)

        # Theme patchen — dk1 (Textfarbe) auf Weiß setzen
        _patch_theme(save_path)

        cleanup_temp_images(images)
        no_img = "" if any(images) else " (ohne Bilder — Pexels API-Key eintragen)"
        return f"Praesentation gespeichert: {save_path} ({len(prs.slides)} Folien){no_img}"

    except Exception as e:
        import traceback
        return f"Fehler: {e}\n{traceback.format_exc()}"


def create_presentation_from_text(topic: str, outline: str, save_path: str = None) -> str:
    try:
        slides_data = json.loads(outline)
    except Exception:
        lines = [l.strip() for l in outline.split("\n") if l.strip()]
        slides_data = [{"title": "Einleitung", "subtitle": topic}]
        for line in lines:
            slides_data.append({"title": line, "bullets": []})
    return create_presentation(topic, slides_data, save_path)
